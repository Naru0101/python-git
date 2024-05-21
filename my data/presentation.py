from pptx import Presentation
from pptx.util import Inches

# Створення нової презентації
prs = Presentation()

# Слайд 1: Інтро
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
title.text = "Михайло Кропивницький"
subtitle = slide_title.placeholders[1]
subtitle.text = "Презентація"

# Слайд 2: Життєві дані
slide = prs.slides.add_slide(prs.slide_layouts[1])
title_shape = slide.shapes.title
title_shape.text = "Життєві дані"
content_shape = slide.placeholders[1]
content_shape.text = "Михайло Кропивницький (1840-1910) народився в селі Ново-Ушиця та помер у місті Одеса."

# Слайд 3: Творчість
slide = prs.slides.add_slide(prs.slide_layouts[1])
title_shape = slide.shapes.title
title_shape.text = "Творчість"
content_shape = slide.placeholders[1]
content_shape.text = "Михайло Кропивницький є автором понад 40 оперет."

# Слайд 4: Народна слава
slide = prs.slides.add_slide(prs.slide_layouts[1])
title_shape = slide.shapes.title
title_shape.text = "Народна слава"
content_shape = slide.placeholders[1]
content_shape.text = "Михайло Кропивницький відомий як «Кобзар Оперети»."

# Слайд 5: Визнання
slide = prs.slides.add_slide(prs.slide_layouts[1])
title_shape = slide.shapes.title
title_shape.text = "Визнання"
content_shape = slide.placeholders[1]
content_shape.text = "Михайло Кропивницький отримав численні нагороди та відзнаки."

# Слайд 6: Спадок
slide = prs.slides.add_slide(prs.slide_layouts[1])
title_shape = slide.shapes.title
title_shape.text = "Спадок"
content_shape = slide.placeholders[1]
content_shape.text = "Творчість Михайла Кропивницького продовжує жити у серцях українців."

# Збереження презентації
prs.save("Михайло_Кропивницький.pptx")
